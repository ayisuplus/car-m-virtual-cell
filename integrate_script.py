#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
CAR-M Simulator PPT + Script Integration Script

Maps presentation-script.md content to each slide:
- Slide body: distilled key points (concise, readable)
- Slide notes: full speaker script for that slide
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from copy import deepcopy
import os

# Input / Output paths
INPUT_PPT = 'CAR-M-Simulator-demo-deck-v3.pptx'
OUTPUT_PPT = 'CAR-M-Simulator-demo-deck-v3-integrated.pptx'

# ============================================================
# 1. SCRIPT MAPPING: Full script text per slide
# ============================================================

SCRIPT_MAP = {
    1: """各位老师、各位同学，大家好。

今天我想用九分钟，给大家看一个能跑起来的假设——它叫 CAR-M Simulator，是一个机制启发型的 Web 平台，用来探索 CAR-engineered macrophage，也就是经过 CAR 工程化的巨噬细胞，在肿瘤微环境里的行为动态。

[操作：手指向屏幕标题]""",

    2: """为什么做这个？实体瘤的免疫治疗一直有个头疼的问题。大家应该都熟 CAR-T——它在血液瘤里很漂亮，但一到实体瘤就麻烦：进不去、被压制，抗原还乱，毒性也难管。

CAR-M 厉害就厉害在这：巨噬细胞天然能钻进实体组织，而且它不只会杀伤，还能吞噬、做抗原呈递，甚至重塑微环境。

我们的目标很朴素——做一个可调参、可运行、可比较的交互系统，把「巨噬细胞怎么极化、吞不吞得动肿瘤」这件事摆到屏幕上看。""",

    3: """但这里我要先划一条线，也主动跟各位老师交个底：

我们这个平台，做的是机制趋势的探索，不是临床预测。

它把 CAR 设计、TME 条件和细胞行为连起来，帮我们比较「如果这样设计会怎样」，但它不能替代实验、也不能替代临床判断，最后一定要用实验去校准。

这条边界，后面我还会再提一次。""",

    4: """好，接下来我切到网站，大家先看三个东西就好，不用记按钮。

先给大伙一个「这是活的」的感觉。这是我们用 AI 生成的 DNA 双螺旋 3D 模型，可以转着看。待会儿模拟里那些细胞，就是在这样的微观尺度上跑的。

第一，细胞行为。你看这块——CAR-M、肿瘤细胞、野生型巨噬、还有 CD8+ T 细胞，它们在同一个微环境里实时运动、互相作用，不是动画，是模拟在跑。

第二，模型输入。左边这些——CAR 信号域、抗原、亲和力、CD47/SIRPα 阻断，还有 TME 条件——都会改变模拟的输入。

第三，结果反馈。右边这个 Dashboard 实时追踪四个东西：肿瘤数量、吞噬率、M1/M2 比例，还有 CD8+ T 细胞的活化趋势。

所以等下你们盯住一个点就行：我一动左边参数，右边指标就动。这个「动」，就是今天的核心。""",

    5: """我先离开模拟台，看一眼它背后是什么。

它不是只套了层 UI 皮，背后是四层。我一层一层说：

- 最底下是数据层——scRNA 测序、巨噬细胞图谱，这是原料；
- 往上一层是 AI 层，借鉴 scVI、CellForge 这一类单细胞降维和生成思路，加一个 surrogate 模型。参数一部分来自这些数据，一部分是假设性的，后面要靠实验校准；
- 再往上是模拟层，ABM——也就是基于智能体的建模——加上扩散场、极化方程、吞噬概率；
- 最顶上是应用层，React 加 Three.js，把它们变成你刚才能点的界面。""",

    8: """这里有个我觉得最值得说的小工程点。

极化方程——就是每个细胞每帧都要算它往 M1 还是 M2 极化——如果老老实实解 ODE，一百多个细胞实时跑，每个细胞每帧 847 毫秒，整个画面就卡死了。

所以我们用了一个神经 surrogate，把这一步压到 0.3 毫秒。

你看，我在网站上可以跑一个参数扫描——左边是原来的 ODE 速度，右边是 surrogate 推理速度。快不是为了炫技——快，才让「每一帧给上百个细胞做决策」这件事变成可能。

所以你刚看到画面里细胞动得那么顺，背后就是它每帧替每个细胞算完了一遍。""",

    7: """回到模拟台，下面是今天最想给大家看的部分。我按一条固定路线走，不随机乱点。

第一步，建立共同语境。我先跑一个 baseline——它受 CT-0508 启发，CT-0508 是首个进临床的 HER2 CAR-M 试验，收尾我会提它的真实数据。

现在大家看默认设置：CAR-M 数量、肿瘤细胞数量、TGF-β 水平，还有一个随机种子，都先固定。点运行。画面里细胞开始动，右边指标开始爬——肿瘤数量在掉，吞噬率在升，这是基线趋势。

第二步，证明它能探索。现在我把抗原情境换一下，比如 HER2-low，或者勾上 CD147 的 ECM degradation，或者切到 cold tumor 场景。你看，只是输入变了，右边的趋势线形状就变了——M1/M2 比例、CD8 活化都跟着动。

这里我再把早上那条线拎出来一下：我比的是趋势，不是在预测哪个设计能进临床。哪个掉得快、哪个活化强，是方向性的比较，不是结论。

第三步，做对比。最后把前后两个跑分放一起比。不用我多解释，两条线一摆——哪个设计让肿瘤掉得快、哪个让免疫活化强，一眼就看到。

所以各位不妨想一下：如果是你，会先动哪个参数？没有标准答案，但这正是我们想用这个平台帮大家干的事。""",

    6: """收尾前，我想特别谨慎地说一组数字。

刚才提到 CT-0508，它的早期报告里：14 个病人，Grade≥3 的 CRS 是 0%，44% 疾病稳定，最佳反应中位下降 20%。

但我要强调——这是背景和场景设计的参考，不是我们这个模拟的验证结果。我们没说模型预测了临床。""",

    9: """最后想说一句话：让机制假设先在屏幕上跑起来。

下一步很明确——把更多实验数据接进来做参数校准，让 CAR-M 设计、TME 情境和验证实验形成一个闭环。我们不会去宣称模型已经能预测临床。

我希望它不只是个演示，而是来这边之后能和各位的实验、数据真正接上。

今天最想请各位给的反馈是三件事：机制边界划得对不对、surrogate 这种折中能不能接受、还有后续该接哪些实验来验证。

谢谢大家！""",
}


# ============================================================
# 2. BODY CONTENT: Distilled key points for each slide
# ============================================================

BODY_UPDATES = {
    # Slide 1: Title - keep mostly as-is, minor polish
    1: {
        "title": "CAR-M Simulator",
        "subtitle": "用于探索 CAR-engineered macrophage 在肿瘤微环境中行为动态的机制启发型 Web 平台",
        "tagline": "AI Virtual Cell Platform · Computational Immunology & Systems Biology",
    },

    # Slide 2: Background - concise three columns
    2: {
        "heading": "实体瘤免疫治疗的建模问题",
        "body_main": "CAR-T 在实体瘤中常受限于浸润不足、免疫抑制性 TME、抗原异质性和毒性管理。CAR-M 天然进入实体组织，同时承担吞噬、抗原呈递与微环境重塑。",
        "col1_title": "障碍",
        "col1_body": "抑制性细胞因子、低氧/乳酸、ECM 屏障改变免疫细胞行为。",
        "col2_title": "机会",
        "col2_body": "巨噬细胞极化状态和吞噬概率可被工程化设计影响。",
        "col3_title": "展示目标",
        "col3_body": "把机制假设转成可调参、可运行、可比较的交互系统。",
    },

    # Slide 3: Core Positioning - emphasize boundary
    3: {
        "title": "机制趋势，不是临床预测",
        "body": "平台把 CAR 设计、TME 条件、细胞行为和实时指标连接起来，用于比较机制趋势，而不是替代实验或临床判断。",
    },

    # Slide 4: Demo - three key points
    4: {
        "heading": "Demo 时请盯住三个变化",
        "step1_title": "细胞行为",
        "step1_body": "CAR-M、肿瘤细胞、野生型巨噬和 CD8+ T 细胞在同一 TME 中实时运动、相互作用。",
        "step2_title": "模型输入",
        "step2_body": "CAR 信号域、抗原、亲和力、CD47/SIRPα 阻断和 TME 条件改变模拟输入。",
        "step3_title": "结果反馈",
        "step3_body": "Dashboard 实时追踪肿瘤数量、吞噬率、M1/M2 比例和 CD8+ T 活化趋势。",
        "footer": "切网站：进入 Simulation Workbench，先播放默认场景，再打开右侧 Metrics 面板。",
    },

    # Slide 5: Architecture - four layers
    5: {
        "heading": "四层把生物知识、机器学习和可交互模拟接起来",
        "layer1_title": "Data Layer",
        "layer1_body": "scRNA-seq / TAM atlas / 通路与结构参考。",
        "layer2_title": "AI / ML Layer",
        "layer2_body": "scVI + CellForge 思路 + surrogate。参数部分来自数据，部分为假设性，需实验校准。",
        "layer3_title": "Simulation Layer",
        "layer3_body": "ABM、扩散场、极化 ODE、吞噬概率。",
        "layer4_title": "Application Layer",
        "layer4_body": "React、Canvas / Three.js、Chart.js。",
    },

    # Slide 8: Engineering Highlight (technical)
    8: {
        "heading": "神经 surrogate 让极化模型进入实时 ABM",
        "col1_title": "ODE 太慢",
        "col1_body": "每细胞每帧求解极化 ODE，100+ 细胞实时决策被拖慢。",
        "col2_title": "细胞因子",
        "col2_body": "IFN-γ、IL-4、IL-10、TGF-β 等局部信号输入模型。",
        "col3_title": "M1 / M2",
        "col3_body": "输出极化分数，参与吞噬概率和细胞行为更新。",
        "col4_title": "0.3ms",
        "col4_body": "ODE 约 847ms → surrogate 推理约 0.3ms。快让「每帧上百细胞决策」成为可能。",
        "footer": "切网站：打开 AI 面板，运行 parameter sweep。",
    },

    # Slide 7: Demo Route
    7: {
        "heading": "按基线到调参再到对比的路线走",
        "step1_title": "先建立共同语境",
        "step1_body": "运行 CT-0508-inspired baseline，展示默认细胞组成、CAR-M 数量、肿瘤细胞数量、TGF-β 和随机种子。",
        "step2_title": "再证明可探索性 & 做对比",
        "step2_body": "切到 CAR-M Designer / Preset Scenarios，展示 HER2-low、CD147 ECM degradation、cold tumor 等情境如何改变趋势。打开 Compare 面板，放 baseline vs 调参后，一眼看出哪个设计让肿瘤掉得快、免疫活化强。",
        "footer": "切网站：Scenarios → CAR-M → Metrics → Compare。",
    },

    # Slide 6: Clinical Anchor
    6: {
        "heading": "CT-0508 早期报告：14 个病人，0% Grade≥3 CRS，44% 疾病稳定，最佳反应中位下降 20%",
        "disclaimer": "本项目仅把这些结果作为临床背景和场景设计参考，不是本模拟的验证结果。",
    },

    # Slide 9: Closing
    9: {
        "title": "让机制假设先在屏幕上跑起来",
        "body": "下一步不是宣称模型已预测临床，而是把更多实验数据接入参数校准，让 CAR-M 设计、TME 情境和验证实验形成闭环。",
    },
}


# ============================================================
# 3. PROCESSING FUNCTIONS
# ============================================================

def find_shape_by_text(slide, text_substring, max_results=5):
    """Find shapes whose text contains the given substring."""
    matches = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            if text_substring in shape.text_frame.text:
                matches.append(shape)
                if len(matches) >= max_results:
                    break
    return matches


def set_text_content(shape, new_text, preserve_format=True):
    """Replace all text in a shape while preserving paragraph formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if len(tf.paragraphs) == 0:
        return
    # Clear all paragraphs and recreate
    p = tf.paragraphs[0]
    # Save first run formatting if available
    saved_font = None
    if p.runs:
        saved_font = p.runs[0].font
    p.clear()
    run = p.add_run()
    # Restore font formatting
    if saved_font:
        if saved_font.name:
            run.font.name = saved_font.name
        if saved_font.size is not None:
            run.font.size = saved_font.size
        if saved_font.bold is not None:
            run.font.bold = saved_font.bold
        if saved_font.italic is not None:
            run.font.italic = saved_font.italic
        if saved_font.color.rgb is not None:
            run.font.color.rgb = saved_font.color.rgb
    run.text = new_text
    # Remove extra paragraphs
    for para in tf.paragraphs[1:]:
        para.clear()
        p_elm = para._element
        p_elm.getparent().remove(p_elm)


def set_slide_notes(slide, notes_text):
    """Set or replace notes text for a slide."""
    if not slide.has_notes_slide:
        notes_slide = slide.notes_slide
    else:
        notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = notes_text


def process_slide_1(slide, updates, script):
    """Process title slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if "CAR-M Simulator" in text and len(text) < 30:
            set_text_content(shape, updates["title"])
        elif "用于探索 CAR-engineered" in text:
            set_text_content(shape, updates["subtitle"])
        elif "AI Virtual Cell Platform" in text:
            set_text_content(shape, updates["tagline"])
    set_slide_notes(slide, script)


def process_slide_2(slide, updates, script):
    """Process background slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if "实体瘤免疫治疗的建模问题" in text:
            set_text_content(shape, updates["heading"])
        elif text.startswith("CAR-T 在实体瘤"):
            set_text_content(shape, updates["body_main"])
        elif text == "障碍":
            set_text_content(shape, updates["col1_title"])
        elif "抑制性细胞因子" in text and len(text) < 50:
            set_text_content(shape, updates["col1_body"])
        elif text == "机会":
            set_text_content(shape, updates["col2_title"])
        elif "巨噬细胞的极化状态" in text:
            set_text_content(shape, updates["col2_body"])
        elif text == "展示目标":
            set_text_content(shape, updates["col3_title"])
        elif "把机制假设转成" in text:
            set_text_content(shape, updates["col3_body"])
    set_slide_notes(slide, script)


def process_slide_3(slide, updates, script):
    """Process core positioning slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if "机制趋势" in text and len(text) < 20:
            set_text_content(shape, updates["title"])
        elif text.startswith("平台把 CAR"):
            set_text_content(shape, updates["body"])
    set_slide_notes(slide, script)


def process_slide_4(slide, updates, script):
    """Process website demo slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if "Demo 时请盯住" in text:
            set_text_content(shape, updates["heading"])
        elif text == "细胞行为":
            set_text_content(shape, updates["step1_title"])
        elif "CAR-M、肿瘤细胞、野生型" in text:
            set_text_content(shape, updates["step1_body"])
        elif text == "模型输入":
            set_text_content(shape, updates["step2_title"])
        elif "CAR 信号域、抗原" in text:
            set_text_content(shape, updates["step2_body"])
        elif text == "结果反馈":
            set_text_content(shape, updates["step3_title"])
        elif "Dashboard 追踪" in text:
            set_text_content(shape, updates["step3_body"])
        elif "切网站：进入 Simulation" in text:
            set_text_content(shape, updates["footer"])
    set_slide_notes(slide, script)


def process_slide_5(slide, updates, script):
    """Process architecture slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if "四层把生物知识" in text:
            set_text_content(shape, updates["heading"])
        elif text == "Data Layer":
            set_text_content(shape, updates["layer1_title"])
        elif "scRNA-seq / TAM" in text:
            set_text_content(shape, updates["layer1_body"])
        elif text == "AI / ML Layer":
            set_text_content(shape, updates["layer2_title"])
        elif "scVI + CellForge" in text:
            set_text_content(shape, updates["layer2_body"])
        elif text == "Simulation Layer":
            set_text_content(shape, updates["layer3_title"])
        elif "ABM、扩散场" in text:
            set_text_content(shape, updates["layer3_body"])
        elif text == "Application Layer":
            set_text_content(shape, updates["layer4_title"])
        elif "React、Canvas" in text:
            set_text_content(shape, updates["layer4_body"])
    set_slide_notes(slide, script)


def process_slide_8(slide, updates, script):
    """Process engineering highlight slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if "神经 surrogate" in text:
            set_text_content(shape, updates["heading"])
        elif text == "ODE 太慢":
            set_text_content(shape, updates["col1_title"])
        elif "每细胞每帧求解" in text:
            set_text_content(shape, updates["col1_body"])
        elif text == "细胞因子":
            set_text_content(shape, updates["col2_title"])
        elif "IFN-γ、IL-4" in text:
            set_text_content(shape, updates["col2_body"])
        elif text == "M1 / M2":
            set_text_content(shape, updates["col3_title"])
        elif "输出极化分数" in text:
            set_text_content(shape, updates["col3_body"])
        elif text == "0.3ms":
            set_text_content(shape, updates["col4_title"])
        elif "ODE 约 847ms" in text:
            set_text_content(shape, updates["col4_body"])
        elif "切网站：打开 AI" in text:
            set_text_content(shape, updates["footer"])
    set_slide_notes(slide, script)


def process_slide_7(slide, updates, script):
    """Process demo route slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if "按基线到调参" in text:
            set_text_content(shape, updates["heading"])
        elif text == "先建立共同语境":
            set_text_content(shape, updates["step1_title"])
        elif "运行 CT-0508" in text:
            set_text_content(shape, updates["step1_body"])
        elif text == "再证明可探索性":
            set_text_content(shape, updates["step2_title"])
        elif "切到 CAR-M Designer" in text:
            set_text_content(shape, updates["step2_body"])
        elif "做对比" in text and len(text) < 10:
            # Need to be careful with short text matches
            set_text_content(shape, updates.get("step3_title", "做对比"))
        elif "打开 Compare" in text:
            set_text_content(shape, updates["step3_body"])
        elif "切网站：Scenarios" in text:
            set_text_content(shape, updates["footer"])
    set_slide_notes(slide, script)


def process_slide_6(slide, updates, script):
    """Process clinical anchor slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if "CT-0508 早期报告" in text or "Grade" in text or "CRS" in text:
            if len(text) > 30:
                set_text_content(shape, updates["heading"] + "\n\n" + updates["disclaimer"])
    set_slide_notes(slide, script)


def process_slide_9(slide, updates, script):
    """Process closing slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if "让机制假设先在屏幕上跑起来" in text:
            set_text_content(shape, updates["title"])
        elif text.startswith("下一步不是宣称"):
            set_text_content(shape, updates["body"])
    set_slide_notes(slide, script)


# ============================================================
# 4. MAIN EXECUTION
# ============================================================

def main():
    prs = Presentation(INPUT_PPT)
    print(f"Loaded {INPUT_PPT}: {len(prs.slides)} slides")

    # Process each slide
    for i, slide in enumerate(prs.slides, 1):
        print(f"Processing slide {i}...")
        if i in BODY_UPDATES and i in SCRIPT_MAP:
            if i == 1:
                process_slide_1(slide, BODY_UPDATES[i], SCRIPT_MAP[i])
            elif i == 2:
                process_slide_2(slide, BODY_UPDATES[i], SCRIPT_MAP[i])
            elif i == 3:
                process_slide_3(slide, BODY_UPDATES[i], SCRIPT_MAP[i])
            elif i == 4:
                process_slide_4(slide, BODY_UPDATES[i], SCRIPT_MAP[i])
            elif i == 5:
                process_slide_5(slide, BODY_UPDATES[i], SCRIPT_MAP[i])
            elif i == 8:
                process_slide_8(slide, BODY_UPDATES[i], SCRIPT_MAP[i])
            elif i == 7:
                process_slide_7(slide, BODY_UPDATES[i], SCRIPT_MAP[i])
            elif i == 6:
                process_slide_6(slide, BODY_UPDATES[i], SCRIPT_MAP[i])
            elif i == 9:
                process_slide_9(slide, BODY_UPDATES[i], SCRIPT_MAP[i])
        else:
            print(f"  Skipped (no mapping)")

    # Save output
    prs.save(OUTPUT_PPT)
    print(f"\nSaved integrated presentation to: {OUTPUT_PPT}")
    print(f"Full speaker scripts added to notes of all {len(prs.slides)} slides.")


if __name__ == '__main__':
    main()
