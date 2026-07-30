#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Finalize the CAR-M Simulator deck for delivery.

Base: CAR-M-Simulator-demo-deck-v3-optimized.pptx (richest visual version).
Changes (surgical, content-preserving):
  1. Fix the ONE unsupported performance claim on Slide 8.
     Old: big number "0.3ms"; caption "网站里可对比 ODE 约 847ms 与 surrogate 推理约 0.3ms。"
     New: measured benchmark — surrogate ~1.5µs vs ODE ~12µs, ~8x speedup.
  2. Embed the full speaker script into each slide's notes (from presentation-script.md),
     so the deck is self-contained on stage.
Everything else (layout, colors, other copy) is preserved as-is.
"""
from pptx import Presentation

SRC = 'CAR-M-Simulator-demo-deck-v3-optimized.pptx'
OUT = 'CAR-M-Simulator-demo-deck-v3-final.pptx'

prs = Presentation(SRC)

# ---- 1. Fix Slide 8 performance numbers ----------------------------------
slide8 = prs.slides[7]
for sh in slide8.shapes:
    if not sh.has_text_frame:
        continue
    txt = sh.text_frame.text.strip()
    # big-number card
    if txt == '0.3ms':
        # keep styling of the first run, replace text
        para = sh.text_frame.paragraphs[0]
        if para.runs:
            para.runs[0].text = '~8x'
            for r in para.runs[1:]:
                r.text = ''
        else:
            para.text = '~8x'
    # caption under it
    elif '847ms' in txt or '0.3ms' in txt:
        new_caption = ('实测：surrogate 推理约 1.5µs / 细胞，等价 ODE 求解约 12µs，'
                       '约 8× 提速；N≈55 细胞时每帧极化 < 0.1ms，60fps 下余量充足。')
        para = sh.text_frame.paragraphs[0]
        if para.runs:
            para.runs[0].text = new_caption
            for r in para.runs[1:]:
                r.text = ''
        else:
            para.text = new_caption

# ---- 2. Full speaker script per slide (from presentation-script.md v3) ----
NOTES = {
1: """【开场 ~1.5′｜PPT1→2→3】各位老师、各位同学，大家好。
今天我想用九分钟，给大家看一个能跑起来的假设——它叫 CAR-M Simulator，一个机制启发型 Web 平台，用来探索经 CAR 工程化的巨噬细胞在肿瘤微环境里的行为动态。[手指屏幕标题]""",
2: """为什么做？实体瘤免疫治疗一直头疼。CAR-T 在血液瘤漂亮，一到实体瘤就麻烦：进不去、被压制、抗原乱、毒性难管。
CAR-M 强在：天然钻进实体组织，且不只杀伤，还能吞噬、抗原呈递、重塑微环境。
目标很朴素——做一个可调参、可运行、可比较的交互系统。""",
3: """[放慢，强调]先划一条线，也跟各位交底：这个平台做的是机制趋势探索，不是临床预测。
它把 CAR 设计、TME 条件、细胞行为连起来比较「如果这样设计会怎样」，但不替代实验、不替代临床判断，最后一定要实验校准。这条边界后面还会再提。""",
4: """[切网站]大家看三个东西就好，不用记按钮。
先转一下 AI 生成的 DNA 双螺旋 3D 模型，给个「这是活的」感觉。路线：3D → Simulation → Metrics。
①细胞行为：CAR-M/肿瘤/野生型巨噬/CD8+ T 在同一 TME 实时运动——不是动画，是模拟在跑。
②模型输入：CAR 信号域、抗原、亲和力、CD47/SIRPα 阻断、TME 条件。
③结果反馈：右侧 Dashboard 实时追踪肿瘤数量、吞噬率、M1/M2、CD8 活化。
一句话:我一动左边，右边就动。""",
5: """[离开模拟台]它不只是套层 UI 皮，背后四层：
①数据层——scRNA-seq / TAM atlas，原料;
②AI 层——借鉴 scVI、CellForge 思路 + surrogate，参数部分来自数据、部分假设性,需校准;
③模拟层——ABM + 扩散场 + 极化方程 + 吞噬概率;
④应用层——React + Three.js，变成你能点的界面。""",
6: """[谨慎报数]CT-0508 早期报告：14 个病人，Grade≥3 CRS 为 0%，44% 疾病稳定，最佳反应中位下降 20%。
强调——这是背景和场景设计参考，不是我们模拟的验证结果。我们没说模型预测了临床。""",
7: """[回模拟台]按固定路线走，不乱点。Scenarios → CAR-M → Metrics → Compare。
A 建立共同语境:跑 CT-0508-inspired baseline，固定 CAR-M 数量、肿瘤数量、TGF-β、随机种子，点 Run。
B 证明可探索:换 HER2-low / CD147 ECM degradation / cold tumor，只改输入,看趋势线变化。
再拎一次边界:比的是趋势，不是预测哪个能进临床。
C 对比:Compare 面板把 baseline vs 调参后并排一摆，一眼看到差异。""",
8: """[最想说的工程点]极化方程每个细胞每帧都要算它往 M1 还是 M2 极化。老实解 ODE，上百细胞实时跑会卡。
我们用神经 surrogate 替代这一步。实测:surrogate 推理约 1.5µs/细胞，等价 ODE 求解约 12µs，约 8× 提速;N≈55 细胞时每帧极化 <0.1ms，60fps 余量充足。
[打开 AI 面板跑 parameter sweep]快不是炫技——快才让「每帧给上百细胞做决策」变可能。画面动得顺，就是因为每帧替每个细胞算完了。
(注:早期材料写的 847ms/0.3ms 无源码依据，已弃用,改用可复现的实测值。)""",
9: """[收尾]最后一句:让机制假设先在屏幕上跑起来。
下一步不是宣称已预测临床，而是把更多实验数据接入参数校准,让 CAR-M 设计、TME 情境、验证实验形成闭环。
今天最想请各位反馈三件事:机制边界划得对不对、surrogate 折中能不能接受、后续该接哪些实验来验证。谢谢大家!""",
}

for i, slide in enumerate(prs.slides, 1):
    if i in NOTES:
        slide.notes_slide.notes_text_frame.text = NOTES[i]

prs.save(OUT)
print('Saved:', OUT)

# verify slide 8
p2 = Presentation(OUT)
print('--- Slide 8 after fix ---')
for sh in p2.slides[7].shapes:
    if sh.has_text_frame and sh.text_frame.text.strip():
        t = sh.text_frame.text.strip()
        if '8x' in t or 'µs' in t or 'surrogate' in t or 'ODE' in t:
            print(repr(t)[:120])
print('--- notes lengths ---', {i: len(prs.slides[i-1].notes_slide.notes_text_frame.text) for i in NOTES})
