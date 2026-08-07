#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Build CAR-M Simulator demo deck v4 — expanded 14-slide version.
Embeds real website screenshots + RunningHub scientific illustrations.
Self-contained: can be understood without live demo.
"""
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Paths
BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHOTS = os.path.join(BASE, 'assets-poster', 'screenshots')
ART = os.path.join(BASE, 'assets-poster')
OUT = os.path.join(BASE, 'CAR-M-Simulator-demo-deck-v4.pptx')

# Colors
NAVY = RGBColor(0x0B, 0x14, 0x26)
DARK_BLUE = RGBColor(0x10, 0x1E, 0x3A)
TEAL = RGBColor(0x00, 0xD4, 0xAA)
CYAN = RGBColor(0x00, 0xC8, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xDD, 0xEE)
ACCENT = RGBColor(0x4E, 0xC9, 0xB0)
ORANGE = RGBColor(0xFF, 0x8C, 0x42)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color=NAVY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_name='Calibri'):
    """lines: list of (text, size, color, bold, align)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(6)
    return txBox

def add_image_fit(slide, path, left, top, max_w, max_h):
    """Add image maintaining aspect ratio within bounds."""
    if not os.path.exists(path):
        print(f"  [WARN] missing: {path}")
        return None
    from PIL import Image
    im = Image.open(path)
    iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w = int(iw * ratio)
    h = int(ih * ratio)
    cx = left + (max_w - w) // 2
    cy = top + (max_h - h) // 2
    return slide.shapes.add_picture(path, cx, cy, w, h)

def add_card(slide, left, top, width, height, fill_color=DARK_BLUE, radius=Inches(0.15)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0x1A, 0x3A, 0x5C)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape

def set_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

W = prs.slide_width
H = prs.slide_height

# ============================================================
# SLIDE 1: COVER
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_shape_bg(s, NAVY)
# Hero image as background (right side)
hero_path = os.path.join(ART, 'hero-car-m.png')
if os.path.exists(hero_path):
    s.shapes.add_picture(hero_path, Inches(5.5), Inches(0), Inches(7.833), Inches(7.5))
# Gradient overlay left
overlay = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(7), H)
overlay.fill.solid()
overlay.fill.fore_color.rgb = NAVY
overlay.line.fill.background()
# Title
add_multi_text(s, Inches(0.8), Inches(1.5), Inches(6), Inches(5), [
    ("CAR-M Simulator", 44, TEAL, True, PP_ALIGN.LEFT),
    ("", 12, WHITE, False, PP_ALIGN.LEFT),
    ("A Mechanism-Inspired, Real-Time", 22, WHITE, False, PP_ALIGN.LEFT),
    ("Agent-Based Platform for Exploring", 22, WHITE, False, PP_ALIGN.LEFT),
    ("CAR-Engineered Macrophage Behavior", 22, WHITE, False, PP_ALIGN.LEFT),
    ("in the Tumor Microenvironment", 22, WHITE, False, PP_ALIGN.LEFT),
    ("", 12, WHITE, False, PP_ALIGN.LEFT),
    ("Cambridge Visiting Research Program  |  2025", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
])
set_notes(s, """Cover slide. Stand to the side of the screen, greet audience.
Today I want to show you a hypothesis that runs - CAR-M Simulator, a mechanism-inspired web platform for exploring CAR-engineered macrophage behavior in the tumor microenvironment.""")

# ============================================================
# SLIDE 2: BACKGROUND & MOTIVATION
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(s, NAVY)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Why CAR-M?", 36, TEAL, True)
# Three columns
col_data = [
    ("CAR-T Limitations\nin Solid Tumors", [
        "Poor tissue infiltration",
        "Immunosuppressive TME",
        "Antigen heterogeneity",
        "Systemic toxicity (CRS)",
    ]),
    ("CAR-M Advantages", [
        "Native solid tissue motility",
        "Phagocytosis + antigen presentation",
        "TME remodeling capacity",
        "CT-0508: 0% CRS >= G3",
    ]),
    ("Our Goal", [
        "A configurable, runnable,",
        "comparable interactive system",
        "to visualize macrophage",
        "polarization & phagocytosis",
    ]),
]
for i, (title, bullets) in enumerate(col_data):
    x = Inches(0.8 + i * 4.1)
    add_card(s, x, Inches(1.5), Inches(3.7), Inches(5.2))
    add_text_box(s, x + Inches(0.3), Inches(1.7), Inches(3.1), Inches(1),
                 title, 16, CYAN, True)
    lines = [(f"  {b}", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT) for b in bullets]
    add_multi_text(s, x + Inches(0.3), Inches(2.8), Inches(3.1), Inches(3.5), lines)
set_notes(s, """CAR-T works beautifully in blood cancers, but solid tumors are a different story: can't get in, get suppressed, antigen is messy, toxicity is hard to manage.
CAR-M is powerful because macrophages naturally infiltrate solid tissue, and they don't just kill - they phagocytose, present antigens, and remodel the microenvironment.
Our goal is simple - make a configurable, runnable, comparable interactive system.""")

# ============================================================
# SLIDE 3: CORE POSITIONING (Boundary)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(s, NAVY)
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Scope & Boundary", 36, ORANGE, True)
add_card(s, Inches(1.5), Inches(1.8), Inches(10.3), Inches(4.5))
add_multi_text(s, Inches(2), Inches(2.2), Inches(9.3), Inches(4), [
    ("Mechanism-Inspired Trend Exploration", 24, TEAL, True, PP_ALIGN.CENTER),
    ("NOT Clinical Prediction", 24, ORANGE, True, PP_ALIGN.CENTER),
    ("", 14, WHITE, False, PP_ALIGN.LEFT),
    ("We connect CAR design, TME conditions, and cell behavior to compare", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ('"what if we design it this way?"', 16, WHITE, True, PP_ALIGN.CENTER),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("It does NOT replace experiments or clinical judgment.", 15, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("Final validation must come from wet-lab calibration.", 15, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("Parameters: partly literature/scRNA atlas, partly plausibility-based assumptions.", 13, RGBColor(0x88,0xAA,0xCC), False, PP_ALIGN.CENTER),
])
set_notes(s, """Let me draw a line upfront: this platform explores mechanistic trends, NOT clinical prediction. It connects CAR design, TME conditions, and cell behavior to compare 'what if we design it this way', but it cannot replace experiments or clinical judgment. Final validation must come from experiments. I'll repeat this boundary later.""")

# ============================================================
# SLIDE 4: PLATFORM OVERVIEW (website hero screenshot)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(s, NAVY)
add_text_box(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Platform Overview", 32, TEAL, True)
add_text_box(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
             "Full-stack web application: React + Three.js + Chart.js | ~14,000 lines TypeScript | Client-side only", 14, LIGHT_GRAY, False)
shot = os.path.join(SHOTS, '01-hero.png')
add_image_fit(s, shot, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.7))
set_notes(s, """Switch to website. Everyone just look at three things, no need to memorize buttons.
First a 'this is alive' feeling. This is our AI-generated DNA double helix 3D model, you can rotate it. The cells in the simulation run at this microscopic scale.
Route: 3D -> Simulation -> Metrics.""")

# ============================================================
# SLIDE 5: 3D CELL MODELS
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(s, NAVY)
add_text_box(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Interactive 3D Cell Viewer", 32, TEAL, True)
add_text_box(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
             "AI-generated 3D models (RunningHub + Meshy) | Rotate, zoom, explore cell morphology", 14, LIGHT_GRAY, False)
shot = os.path.join(SHOTS, '02-3d-viewer.png')
add_image_fit(s, shot, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.7))
set_notes(s, """Spin the 3D model for about 10 seconds to give the audience a 'this is alive' feeling. These are AI-generated models of CAR-macrophages, tumor cells, and DNA structures.""")

# ============================================================
# SLIDE 6: SIMULATION WORKBENCH (running)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(s, NAVY)
add_text_box(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Real-Time Agent-Based Simulation", 32, TEAL, True)
add_text_box(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
             "4 cell types | 10-factor reaction-diffusion field | Deterministic (seed=20250706) | 60 fps", 14, LIGHT_GRAY, False)
shot = os.path.join(SHOTS, '03-simulation-running.png')
add_image_fit(s, shot, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.7))
set_notes(s, """Three things to watch:
1. Cell behavior - CAR-M, tumor, WT macrophage, CD8+ T cells all moving in real-time TME. Not animation, simulation is running.
2. Model inputs - CAR signaling domain, antigen, affinity, CD47/SIRPa blockade, TME conditions.
3. Result feedback - right dashboard tracks tumor count, phagocytosis rate, M1/M2 ratio, CD8 activation.
When I move left parameters, right metrics move. That 'movement' is today's core.""")

# ============================================================
# SLIDE 7: CAR-M DESIGNER & SCENARIOS
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(s, NAVY)
add_text_box(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "CAR-M Designer & Preset Scenarios", 32, TEAL, True)
add_text_box(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
             "Configure signaling domain, antigen, affinity, checkpoint blockade | CT-0508-inspired baseline", 14, LIGHT_GRAY, False)
shot = os.path.join(SHOTS, '07-carm-designer.png')
add_image_fit(s, shot, Inches(0.3), Inches(1.5), Inches(6.2), Inches(5.7))
shot2 = os.path.join(SHOTS, '06-scenarios.png')
add_image_fit(s, shot2, Inches(6.8), Inches(1.5), Inches(6.2), Inches(5.7))
set_notes(s, """Demo route:
Step 1: Run CT-0508-inspired baseline. Fix CAR-M count, tumor count, TGF-beta, random seed. Click Run.
Step 2: Switch to HER2-low or enable CD147 ECM degradation or cold tumor scenario. Only input changes, right trend lines change shape.
Remember: we compare trends, not predict which design reaches clinic.""")

# ============================================================
# SLIDE 8: COMPARE PANEL
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(s, NAVY)
add_text_box(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Side-by-Side Comparison", 32, TEAL, True)
add_text_box(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
             "Deterministic engine ensures reproducible A/B comparison | Baseline vs modified construct", 14, LIGHT_GRAY, False)
shot = os.path.join(SHOTS, '08-compare.png')
add_image_fit(s, shot, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.7))
set_notes(s, """Put two runs side by side in Compare panel. No need to explain much - two lines side by side, which design makes tumor drop faster, which makes immune activation stronger, immediately visible.
Ask audience: if it were you, which parameter would you change first? No standard answer, but that's exactly what we want this platform to help with.""")

# ============================================================
# SLIDE 9: SYSTEM ARCHITECTURE (four layers illustration)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(s, NAVY)
add_text_box(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Four-Layer Architecture", 32, TEAL, True)
# Left: illustration
art_path = os.path.join(ART, 'four-layers.png')
add_image_fit(s, art_path, Inches(0.3), Inches(1.2), Inches(6), Inches(6))
# Right: layer descriptions
layers = [
    ("Data Layer", "scRNA-seq / TAM atlas | Literature priors | Calibratable parameters"),
    ("AI Layer", "Neural surrogate (MLP 6-32-32-3) | Inspired by scVI / CellForge"),
    ("Simulation Layer", "ABM + 10-factor reaction-diffusion | Checkpoint-aware phagocytosis | Polarization dynamics"),
    ("Application Layer", "React + TypeScript + Three.js + Chart.js | ~14K lines | Client-side only"),
]
for i, (title, desc) in enumerate(layers):
    y = Inches(1.5 + i * 1.4)
    add_card(s, Inches(6.8), y, Inches(6), Inches(1.2))
    add_text_box(s, Inches(7.1), y + Inches(0.1), Inches(5.5), Inches(0.5),
                 title, 16, CYAN, True)
    add_text_box(s, Inches(7.1), y + Inches(0.55), Inches(5.5), Inches(0.6),
                 desc, 12, LIGHT_GRAY, False)
set_notes(s, """It's not just a UI skin. Four layers:
Data layer - scRNA-seq, macrophage atlas, raw material.
AI layer - borrows scVI, CellForge approach + surrogate. Parameters partly from data, partly hypothetical, needs calibration.
Simulation layer - ABM + diffusion field + polarization equations + phagocytosis probability.
Application layer - React + Three.js, turns them into the interface you can click.""")

# ============================================================
# SLIDE 10: TME & MECHANISTIC MODELS
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(s, NAVY)
add_text_box(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Tumor Microenvironment Model", 32, TEAL, True)
# Left: TME illustration
tme_path = os.path.join(ART, 'tme-environment.png')
add_image_fit(s, tme_path, Inches(0.3), Inches(1.2), Inches(6), Inches(6))
# Right: key mechanisms
add_card(s, Inches(6.8), Inches(1.3), Inches(6), Inches(5.5))
add_multi_text(s, Inches(7.1), Inches(1.5), Inches(5.5), Inches(5), [
    ("10-Factor Reaction-Diffusion Field", 16, CYAN, True, PP_ALIGN.LEFT),
    ("O2, Lactate, TGF-b, IFN-g, IL-4, IL-10, VEGF, CXCL9, SPP1, ECM", 12, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Checkpoint-Aware Phagocytosis", 16, CYAN, True, PP_ALIGN.LEFT),
    ("p = base x CD47_mod x CD24_mod x M1_bonus x energy", 12, ACCENT, False, PP_ALIGN.LEFT),
    ("CD3z / FcRg / CD147(ECM) / MerTK(efferocytosis)", 12, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Polarization Dynamics", 16, CYAN, True, PP_ALIGN.LEFT),
    ("Local cytokine context -> M1/M2 scores", 12, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("M1: anti-tumor | M2: pro-tumor | MIXED: intermediate", 12, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("CD8+ T-cell Coupling", 16, CYAN, True, PP_ALIGN.LEFT),
    ("IFN-g/CXCL9 recruit -> activation -> clonal expansion", 12, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("Exhaustion under TGF-b/IL-10 load (emergent)", 12, LIGHT_GRAY, False, PP_ALIGN.LEFT),
])
set_notes(s, """The TME model includes 10 factors on a reaction-diffusion grid. Hypoxic core, immunosuppressive cytokine clouds, ECM barriers. The phagocytosis module is the central mechanistic piece - it models the don't-eat-me checkpoint biology.""")

# ============================================================
# SLIDE 11: NEURAL SURROGATE (performance)
# ============================================================
# Read measured benchmark so the deck always agrees with the data
# (scripts/benchmark-surrogate.mjs -> paper/figures/data/benchmark_results.json).
_BENCH_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           '..', 'paper', 'figures', 'data', 'benchmark_results.json')
with open(_BENCH_PATH) as _f:
    _B = json.load(_f)
_SURR_US = _B['surrogate_us']
_ODE_US = _B['ode_us']
_SPEEDUP = _B['speedup']
_FRAME_MS = _B['frame_ms_surrogate_N55']
_HEADROOM = round(16.7 / _FRAME_MS, 0)

s = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(s, NAVY)
add_text_box(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             f"Neural Surrogate: ~{_SPEEDUP:.0f}x Speedup", 32, TEAL, True)
# Left: surrogate illustration
surr_path = os.path.join(ART, 'surrogate-mlp.jpg')
add_image_fit(s, surr_path, Inches(0.3), Inches(1.2), Inches(6.5), Inches(3.5))
# Right: benchmark card
add_card(s, Inches(7), Inches(1.3), Inches(5.8), Inches(3.3))
add_multi_text(s, Inches(7.3), Inches(1.5), Inches(5.2), Inches(3), [
    ("Performance Benchmark (measured, 11 replicates)", 18, CYAN, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    (f"Neural Surrogate:  ~{_SURR_US:.2f} us / cell", 16, ACCENT, True, PP_ALIGN.LEFT),
    (f"ODE (RK4):          ~{_ODE_US:.2f} us / cell", 16, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    (f"Speedup:              ~{_SPEEDUP:.1f}x", 20, TEAL, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    (f"N=55 cells: {_FRAME_MS:.2f} ms/frame (60fps headroom ~{_HEADROOM:.0f}x)", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
])
# Bottom: AI panel screenshot
shot_ai = os.path.join(SHOTS, '09-ai-panel.png')
add_image_fit(s, shot_ai, Inches(0.5), Inches(5), Inches(12.3), Inches(2.3))
set_notes(s, f"""The most worthwhile engineering point. The polarization equation - every cell every frame must compute M1 vs M2. If we honestly solve ODE, 100+ cells real-time will lag every frame. So we use a neural surrogate.
Measured (median of 11 timed replicates): surrogate ~{_SURR_US:.2f}us/cell, representative RK4 ODE ~{_ODE_US:.2f}us, about {_SPEEDUP:.1f}x speedup. N~55 cells per frame polarization {_FRAME_MS:.2f}ms, 60fps margin sufficient.
Open the AI panel parameter sweep on website. Fast is not for showing off - fast makes 'every frame giving 100+ cells a decision' possible.""")

# ============================================================
# SLIDE 12: CLINICAL CONTEXT (CT-0508)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(s, NAVY)
add_text_box(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Clinical Context: CT-0508", 32, ORANGE, True)
add_text_box(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
             "First-in-human HER2 CAR-M trial | Design anchor, NOT validation of our model", 14, LIGHT_GRAY, False)
# Four number cards (all traceable to Reiss et al., Nat Med 2025,
# doi:10.1038/s41591-025-03495-z: 14 patients; no >=G3 CRS; 44% (4/9) HER2 3+
# stable disease as best response; tumor volume reductions in 41% of
# measurable lesions — with NO RECIST objective responses)
nums = [("14", "Patients"), ("0%", "CRS >= G3"), ("44%", "Stable Disease"), ("41%", "Lesions Shrunk")]
for i, (num, label) in enumerate(nums):
    x = Inches(1 + i * 3)
    add_card(s, x, Inches(2), Inches(2.5), Inches(2.5))
    add_text_box(s, x, Inches(2.3), Inches(2.5), Inches(1.2),
                 num, 40, TEAL, True, PP_ALIGN.CENTER)
    add_text_box(s, x, Inches(3.5), Inches(2.5), Inches(0.8),
                 label, 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)
# Warning box
add_card(s, Inches(1), Inches(5), Inches(11.3), Inches(1.8), RGBColor(0x2A, 0x1A, 0x10))
add_multi_text(s, Inches(1.3), Inches(5.2), Inches(10.7), Inches(1.5), [
    ("Important Disclaimer", 16, ORANGE, True, PP_ALIGN.LEFT),
    ("These numbers are background context and scenario design reference.", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("They are NOT validation results of our simulation.", 14, ORANGE, True, PP_ALIGN.LEFT),
    ("We do NOT claim our model predicted clinical outcomes.", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
])
shot_clin = os.path.join(SHOTS, '11-clinical.png')
set_notes(s, """Before wrapping up, I want to carefully report a set of numbers, all from the CT-0508 Phase 1 report (Reiss et al., Nature Medicine 2025, doi:10.1038/s41591-025-03495-z): 14 patients treated; no Grade>=3 CRS or ICANS; 44% (4 of 9) of HER2 3+ patients achieved stable disease as best overall response; tumor volume reductions were noted in 41% of measurable lesions — but importantly, there were NO objective responses per RECIST v1.1, so we do not quote any "tumor shrinkage percentage" as a response number.
Emphasize - this is background and scenario design reference, NOT our simulation's validation result. We did not say the model predicted clinical.""")

# ============================================================
# SLIDE 13: REPRODUCIBILITY & ENGINEERING
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(s, NAVY)
add_text_box(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "Reproducibility & Engineering", 32, TEAL, True)
# Two columns of cards
left_items = [
    ("Deterministic Engine", "Seeded LCG (seed=20250706)\nBit-for-bit reproducible runs\nHonest A/B comparison"),
    ("Client-Side Only", "No server dependency\nStatic bundle ~11s build\nRuns on commodity laptops"),
]
right_items = [
    ("Open & Scripted", "Asset pipeline scripted (RH/Meshy)\nBenchmark ships with code\nConfigurable & shareable"),
    ("Real-Time at Scale", "WebGL/Three.js GPU rendering\nSpatial hashing O(1) neighbor lookup\nResource preloading + fallback"),
]
for i, (title, desc) in enumerate(left_items):
    y = Inches(1.5 + i * 2.8)
    add_card(s, Inches(0.8), y, Inches(5.5), Inches(2.4))
    add_text_box(s, Inches(1.1), y + Inches(0.2), Inches(5), Inches(0.5),
                 title, 16, CYAN, True)
    add_text_box(s, Inches(1.1), y + Inches(0.8), Inches(5), Inches(1.5),
                 desc, 13, LIGHT_GRAY, False)
for i, (title, desc) in enumerate(right_items):
    y = Inches(1.5 + i * 2.8)
    add_card(s, Inches(7), y, Inches(5.5), Inches(2.4))
    add_text_box(s, Inches(7.3), y + Inches(0.2), Inches(5), Inches(0.5),
                 title, 16, CYAN, True)
    add_text_box(s, Inches(7.3), y + Inches(0.8), Inches(5), Inches(1.5),
                 desc, 13, LIGHT_GRAY, False)
set_notes(s, """Engineering highlights: deterministic under fixed seed for honest comparison, entirely client-side with no backend, open and reproducible with scripted asset generation, real-time at scale with spatial hashing and WebGL rendering.""")

# ============================================================
# SLIDE 14: SUMMARY & NEXT STEPS
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(s, NAVY)
# Use hero image as subtle background
if os.path.exists(hero_path):
    pic = s.shapes.add_picture(hero_path, Inches(6), Inches(1), Inches(7.3), Inches(6.5))
    from pptx.oxml.ns import qn
    pic.element.attrib[qn('a:opacity')] = '30000'  # doesn't work in pptx, skip
add_text_box(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Let Mechanistic Hypotheses Run on Screen First", 28, TEAL, True)
add_multi_text(s, Inches(0.8), Inches(1.5), Inches(7), Inches(5), [
    ("What We Built", 20, CYAN, True, PP_ALIGN.LEFT),
    ("A real-time, in-browser ABM platform coupling TME dynamics,", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("neural surrogate, and interactive design workbench.", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("What It Is Good For", 20, CYAN, True, PP_ALIGN.LEFT),
    ("Hypothesis generation | Teaching | Experimental prioritization", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("Next Steps", 20, CYAN, True, PP_ALIGN.LEFT),
    ("1. Data-driven parameter calibration (scRNA-seq / TAM atlas)", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("2. Uncertainty reporting with confidence bands", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("3. Close design -> context -> validation loop", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("Feedback Requested", 20, ORANGE, True, PP_ALIGN.LEFT),
    ("1. Is the mechanism boundary drawn correctly?", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("2. Is the surrogate trade-off acceptable?", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("3. Which experiments should we connect next?", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
])
add_text_box(s, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
             "Thank you!  |  Cambridge Visiting Research Program 2025", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)
set_notes(s, """Last sentence: let mechanistic hypotheses run on screen first.
Next step is not claiming we predicted clinical, but connecting more experimental data for parameter calibration, forming a design-context-validation loop.
Three things I want feedback on today: is the mechanism boundary right, is the surrogate trade-off acceptable, which experiments should we connect next.
Thank you!""")

# ============================================================
# SAVE
# ============================================================
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
